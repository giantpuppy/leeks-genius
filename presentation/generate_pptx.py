from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

# ============= 品牌色 =============
BG = RGBColor(15, 15, 15)      # #0F0F0F
PURPLE = RGBColor(139, 92, 246)  # #8B5CF6
GREEN = RGBColor(52, 211, 153)   # #34D399
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(176, 176, 176)   # #B0B0B0
DARK_GRAY = RGBColor(26, 26, 26) # #1A1A1A

# ============= 尺寸常量 =============
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN_L = Inches(0.7)
MARGIN_R = Inches(0.7)
MARGIN_T = Inches(0.6)
TITLE_H = Inches(0.9)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ============= 工具函数 =============
def set_bg(slide, color=BG):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    # 设置中文字体
    run = p.runs[0]
    run.font._element.set(qn('a:typeface'), font_name)
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=WHITE, line_spacing=1.4):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Microsoft YaHei"
        p.level = 0
        p.space_after = Pt(8)
        p.line_spacing = line_spacing
        run = p.runs[0]
        run.font._element.set(qn('a:typeface'), "Microsoft YaHei")
    return txBox


def add_title(slide, text, color=PURPLE, font_size=40):
    return add_textbox(slide, MARGIN_L, MARGIN_T, SLIDE_W - MARGIN_L - MARGIN_R, TITLE_H,
                       text, font_size=font_size, color=color, bold=True)


def add_subtitle(slide, text, top, color=GREEN, font_size=20, bold=True):
    return add_textbox(slide, MARGIN_L, top, SLIDE_W - MARGIN_L - MARGIN_R, Inches(0.5),
                       text, font_size=font_size, color=color, bold=bold)


def add_placeholder(slide, left, top, width, height, label="截图位置"):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_GRAY
    shape.line.color.rgb = PURPLE
    shape.line.width = Pt(2)

    # 添加占位文字
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(16)
    p.font.color.rgb = PURPLE
    p.font.bold = True
    p.font.name = "Microsoft YaHei"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Inches(height.inches / 2 - 0.3)
    run = p.runs[0]
    run.font._element.set(qn('a:typeface'), "Microsoft YaHei")
    return shape


def add_card(slide, left, top, width, height, title, body, title_color=GREEN, body_color=WHITE):
    # 卡片背景
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_GRAY
    shape.line.color.rgb = PURPLE
    shape.line.width = Pt(1)
    shape.adjustments[0] = 0.08

    # 标题
    add_textbox(slide, left + Inches(0.15), top + Inches(0.12), width - Inches(0.3), Inches(0.35),
                title, font_size=16, color=title_color, bold=True)
    # 正文
    add_textbox(slide, left + Inches(0.15), top + Inches(0.45), width - Inches(0.3), height - Inches(0.6),
                body, font_size=13, color=body_color)


# ============= 第 1 页：封面 =============
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(slide)
add_textbox(slide, Inches(0), Inches(2.2), SLIDE_W, Inches(1.2),
            "排期天菜", font_size=72, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(0), Inches(3.5), SLIDE_W, Inches(0.6),
            "韭菜的自我管理排期剧历", font_size=24, color=PURPLE, bold=True, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(0), Inches(4.2), SLIDE_W, Inches(0.8),
            "让每一张票，都有处安放。", font_size=36, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(0), Inches(5.5), SLIDE_W, Inches(0.5),
            "Paiqi Release", font_size=16, color=GRAY, align=PP_ALIGN.CENTER)

# ============= 第 2 页：排剧期，记场次，存票根 =============
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "排剧期，记场次，存票根")
add_subtitle(slide, "一款面向 杂食党 P 人剧女 的本地排期管理应用", top=Inches(1.5), color=GREEN)
add_textbox(slide, MARGIN_L, Inches(2.0), Inches(6.5), Inches(0.6),
            "把「盘票、记场次、存票根」的碎片流程，变成一个 App。", font_size=18, color=WHITE)
add_textbox(slide, MARGIN_L, Inches(2.8), Inches(6.2), Inches(0.8),
            "盘剧写 repo 快乐，排期不快乐。", font_size=20, color=GRAY, bold=False)

items = [
    "排剧期：可视化对比剧目排期",
    "记场次：已购/想看场次 + 待办提醒",
    "存票根：月底年底数据复盘"
]
add_bullet_list(slide, MARGIN_L, Inches(3.8), Inches(6.2), Inches(2.5), items, font_size=18, color=WHITE)

add_placeholder(slide, Inches(7.3), Inches(1.7), Inches(5.3), Inches(5.0), "App 月历首页 / 开屏页")

# ============= 第 3 页：为谁而做 =============
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "为谁而做")
add_subtitle(slide, "这不是所有人的日历，这是 剧女 的日历", top=Inches(1.45), color=GREEN)

items = [
    "重度戏剧 / 音乐剧爱好者",
    "一年看几十场，杂食党，多剧并行",
    "微信群、大麦、小红书、原生日历四头跑",
    "P 人，不爱整理，但被迫整理"
]
add_bullet_list(slide, MARGIN_L, Inches(2.1), Inches(6.2), Inches(2.4), items, font_size=18, color=WHITE)
add_textbox(slide, MARGIN_L, Inches(4.7), Inches(6.2), Inches(1.2),
            "她们不缺工具，缺的是一套「懂剧场信息密度」的排期工具。", font_size=18, color=GRAY)

add_placeholder(slide, Inches(7.3), Inches(1.6), Inches(5.3), Inches(5.2), "用户场景拼贴")

# ============= 第 4 页：产品解决的痛点与通用工具的局限 =============
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "产品解决的痛点")

# 痛点卡片
card_w = Inches(3.7)
card_h = Inches(2.3)
gap = Inches(0.3)
start_x = MARGIN_L
y = Inches(1.6)

add_card(slide, start_x, y, card_w, card_h,
         "📱 信息散：四处存，找不到",
         "群聊、相册、日历、备忘录各存一点；宣排期截图放相册吃灰，想查卡司时翻半天。")
add_card(slide, start_x + card_w + gap, y, card_w, card_h,
         "⏰ 易出错：午晚场、物料、取票全怕忘",
         "午场 14:00 还是 14:30 常搞混；换物料、帮人取票、面交、买周边、领鸡蛋规则地点全靠临场记忆。")
add_card(slide, start_x + (card_w + gap) * 2, y, card_w, card_h,
         "📊 难复盘：看完就忘，年底一笔糊涂账",
         "票根散落、repo 懒得上传，一年到底看了多少场、花了多少钱、追了哪些演员，没有清晰记录。")

# 通用工具局限
add_subtitle(slide, "通用工具为什么治不好：", top=Inches(4.15), color=GREEN)

# 表格
rows = [
    ("原生日历", "只能记时间，塞不下卡司 / 票版 / 物料"),
    ("Excel / Notion", "太重，P 人坚持不下来"),
    ("大麦 / 猫眼", "只关心「已购票」，不关心「想看」和「规划」"),
    ("小红书收藏", "信息滞后，不好检索")
]

table_left = MARGIN_L
table_top = Inches(4.7)
table_w = SLIDE_W - MARGIN_L - MARGIN_R
table_h = Inches(2.3)
rows_count = len(rows) + 1
cols_count = 2
row_h = table_h / rows_count

# 表头背景
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, table_left, table_top, table_w, row_h)
shape.fill.solid()
shape.fill.fore_color.rgb = PURPLE
shape.line.color.rgb = PURPLE
add_textbox(slide, table_left + Inches(0.1), table_top + Inches(0.05), Inches(2.5), row_h, "工具", font_size=14, color=WHITE, bold=True)
add_textbox(slide, table_left + Inches(2.7), table_top + Inches(0.05), Inches(9.0), row_h, "问题", font_size=14, color=WHITE, bold=True)

for i, (tool, issue) in enumerate(rows):
    y = table_top + row_h * (i + 1)
    # 行背景
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, table_left, y, table_w, row_h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_GRAY if i % 2 == 0 else BG
    shape.line.color.rgb = RGBColor(51, 51, 51)
    add_textbox(slide, table_left + Inches(0.1), y + Inches(0.05), Inches(2.5), row_h, tool, font_size=13, color=WHITE, bold=True)
    add_textbox(slide, table_left + Inches(2.7), y + Inches(0.05), Inches(9.0), row_h, issue, font_size=13, color=GRAY)

# ============= 第 5 页：产品定位 =============
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "产品定位：排 · 记 · 存")
add_textbox(slide, MARGIN_L, Inches(1.4), SLIDE_W - MARGIN_L - MARGIN_R, Inches(0.5),
            "围绕看剧排期的核心流程，做一款简洁美观的本地工具。", font_size=18, color=WHITE)

card_w = Inches(3.9)
card_h = Inches(4.3)
gap = Inches(0.25)
start_x = MARGIN_L
y = Inches(2.0)

add_card(slide, start_x, y, card_w, card_h,
         "🟣 排：横向纵向对比排期流",
         "把剧目宣排期表变成可视化排期流。\n\n支持横向纵向对比场次，一眼看清哪天有哪些剧、哪些卡司，想看的场次手动 mark 后自动同步到月历。")
add_card(slide, start_x + card_w + gap, y, card_w, card_h,
         "🟢 记：保存 + 提醒",
         "已购票场次集中管理，可添加提醒和备注。\n\n换物料、帮人取票、面交、买周边、领鸡蛋规则地点，一场场列清楚，不再临场失忆。")
add_card(slide, start_x + (card_w + gap) * 2, y, card_w, card_h,
         "🟣 存：收录数据，一键复盘",
         "看过的剧目次数、演员次数、花费金额全部变成可视化图表。\n\n一键生成年度看剧报告，让每一张票都有处安放。")

# ============= 第 6 页：排 =============
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "排：把排期表变成可视化排期流")
items = [
    "剧目宣排期一键查看",
    "横向纵向对比场次与卡司",
    "想看的场次手动 mark",
    "自动同步到月历首页"
]
add_bullet_list(slide, MARGIN_L, Inches(1.6), Inches(6.2), Inches(2.5), items, font_size=20, color=WHITE)
add_textbox(slide, MARGIN_L, Inches(4.3), Inches(6.2), Inches(1.0),
            "不再翻相册、不再搜大麦，排期流上一眼盘清楚。", font_size=18, color=GRAY)
add_placeholder(slide, Inches(7.3), Inches(1.5), Inches(5.3), Inches(5.2), "排期板双密度视图")

# ============= 第 7 页：记 =============
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "记：保存 + 提醒，不再临场失忆")
items = [
    "已购票场次集中管理",
    "添加提醒和备注事项",
    "换物料、帮人取票、面交、买周边",
    "领鸡蛋规则地点一一记录"
]
add_bullet_list(slide, MARGIN_L, Inches(1.6), Inches(6.2), Inches(2.5), items, font_size=20, color=WHITE)
add_textbox(slide, MARGIN_L, Inches(4.3), Inches(6.2), Inches(1.0),
            "剧场门口不再手忙脚乱，打开 App 就知道今天要干嘛。", font_size=18, color=GRAY)
add_placeholder(slide, Inches(7.3), Inches(1.5), Inches(5.3), Inches(5.2), "详情页待办清单")

# ============= 第 8 页：存 =============
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "存：月底年底一键复盘")
items = [
    "看过的剧目次数统计",
    "看过的演员次数统计",
    "花费金额可视化",
    "一键生成年度看剧报告"
]
add_bullet_list(slide, MARGIN_L, Inches(1.6), Inches(6.2), Inches(2.5), items, font_size=20, color=WHITE)
add_textbox(slide, MARGIN_L, Inches(4.3), Inches(6.2), Inches(1.0),
            "年底发小红书年度总结，素材直接从这里拿。", font_size=18, color=GRAY)
add_placeholder(slide, Inches(7.3), Inches(1.5), Inches(5.3), Inches(5.2), "个人中心四张图表")

# ============= 第 9 页：视觉设计 =============
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "视觉设计：像素风 + 星之果实紫绿")
items = [
    "黑底 #0F0F0F，护眼不刺眼",
    "紫 #8B5CF6 + 绿 #34D399，星之果实感",
    "像素字体与图标，致敬剧场票根 / 街机复古感",
    "Logo 字谜：「非」+ 横线 +「菜」= 韭，其余为紫"
]
add_bullet_list(slide, MARGIN_L, Inches(1.6), Inches(6.2), Inches(3.0), items, font_size=18, color=WHITE)

# 色块展示
colors = [
    ("#0F0F0F", "背景", DARK_GRAY),
    ("#8B5CF6", "主色紫", PURPLE),
    ("#34D399", "强调绿", GREEN),
]
for i, (hex_val, name, color) in enumerate(colors):
    x = MARGIN_L + i * Inches(1.2)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(4.8), Inches(0.9), Inches(0.9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = WHITE
    add_textbox(slide, x, Inches(5.8), Inches(1.0), Inches(0.4), f"{name}\n{hex_val}", font_size=11, color=GRAY, align=PP_ALIGN.CENTER)

add_placeholder(slide, Inches(7.3), Inches(1.5), Inches(5.3), Inches(5.2), "Logo 拆解图 + 色板")

# ============= 第 10 页：总结 =============
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "总结：我们做对了什么")

card_w = Inches(3.9)
card_h = Inches(3.6)
gap = Inches(0.25)
start_x = MARGIN_L
y = Inches(1.8)

add_card(slide, start_x, y, card_w, card_h,
         "从真实场景出发",
         "不是凭空造工具，而是把自己作为目标用户，把每天重复的排期动作提炼成「排 · 记 · 存」。")
add_card(slide, start_x + card_w + gap, y, card_w, card_h,
         "聚焦核心流程",
         "先做 MVP 验证最小闭环：排期可视化、场次记录、数据复盘，不堆功能。")
add_card(slide, start_x + (card_w + gap) * 2, y, card_w, card_h,
         "视觉差异化",
         "像素风 + 紫绿配色 + Logo 字谜，让「排期天菜」在工具类 App 中有记忆点。")

add_textbox(slide, MARGIN_L, Inches(5.7), SLIDE_W - MARGIN_L - MARGIN_R, Inches(0.8),
            "接下来：OCR 识别排期表、提醒通知、观演记录与社交分享，持续打磨。", font_size=18, color=GRAY)

# ============= 第 11 页：结尾 =============
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_textbox(slide, Inches(0), Inches(2.2), SLIDE_W, Inches(1.2),
            "排期天菜", font_size=72, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(0), Inches(3.5), SLIDE_W, Inches(0.8),
            "让每一张票，都有处安放。", font_size=36, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(0), Inches(4.5), SLIDE_W, Inches(0.6),
            "愿 vibe coding 赐福我的戏梦人生", font_size=24, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(0), Inches(5.5), SLIDE_W, Inches(0.5),
            "Paiqi · 排期天菜", font_size=16, color=GRAY, align=PP_ALIGN.CENTER)

# ============= 保存 =============
prs.save('presentation-styled.pptx')
print("已生成 presentation-styled.pptx")
